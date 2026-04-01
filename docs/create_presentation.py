# -*- coding: utf-8 -*-
"""
Generate PowerPoint presentation: KI-gestützter Feature-Workflow in ReadyStackGo
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Brand colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
SOFT_BG = RGBColor(0xF5, 0xF7, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

MAX_X = 13.0  # safe right margin


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=DARK_GRAY, line_spacing=1.3, font_name="Segoe UI"):
    """lines is list of (text, bold, color_override) or plain str"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, bold, col = item, False, color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            col = item[2] if len(item) > 2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = col
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
    return tf


def add_card(slide, left, top, width, height, title, body_lines,
             title_color=ACCENT_BLUE, bg_color=WHITE):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.space_after = Pt(6)
    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Segoe UI"
        p.space_after = Pt(3)
    return shape


def add_flow_boxes(slide, items, top, box_w, box_h, start_x, gap, arrow_w):
    """Generic horizontal flow with arrow shapes between boxes.
    items: list of (label, sublabel, color)"""
    x = start_x
    for i, (label, sublabel, color) in enumerate(items):
        shape = slide.shapes.add_shape(1, x, top, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_color = WHITE if color != LIGHT_GRAY else DARK_GRAY
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(13) if len(items) <= 5 else Pt(11)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        if sublabel:
            p2 = tf.add_paragraph()
            p2.text = sublabel
            p2.font.size = Pt(10) if len(items) <= 5 else Pt(9)
            p2.font.color.rgb = text_color
            p2.font.name = "Segoe UI"
            p2.alignment = PP_ALIGN.CENTER
        x += box_w + gap
        if i < len(items) - 1:
            arrow = slide.shapes.add_shape(
                13, x, top + box_h // 2 - Pt(6), arrow_w, Pt(12))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LIGHT_GRAY
            arrow.line.fill.background()
            x += arrow_w + gap


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            "KI-gestützter Feature-Workflow", 44, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
            "Von der Idee zum fertigen Feature mit Claude Code", 24, LIGHT_GRAY,
            False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
            "ReadyStackGo  |  2026", 18, ACCENT_BLUE, False, PP_ALIGN.CENTER)

line = slide.shapes.add_shape(1, Inches(4.5), Inches(3.85), Inches(4.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# ============================================================
# SLIDE 2: Agenda
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Agenda", 36, DARK_BG, True)

agenda_items = [
    ("1.", "Überblick: Der KI-gestützte Workflow"),
    ("2.", "Phase 1: Spezifikation erstellen"),
    ("3.", "Phase 2: Feature planen  —  /plan-feature"),
    ("4.", "Phase 3: Feature implementieren  —  /implement-feature"),
    ("5.", "Phase 4: Feature dokumentieren  —  /document-feature"),
    ("6.", "Weitere Skills: /fix-bug, /report-bug, /docker-dev"),
    ("7.", "GitHub-Integration: Issues, Board & Releases"),
    ("8.", "Zusammenfassung & Workflow-Übersicht"),
]

for i, (num, text) in enumerate(agenda_items):
    y = Inches(1.5) + Inches(i * 0.65)
    add_textbox(slide, Inches(1.2), y, Inches(0.6), Inches(0.5), num, 20,
                ACCENT_BLUE, True)
    add_textbox(slide, Inches(1.8), y, Inches(9), Inches(0.5), text, 20,
                DARK_GRAY, False)

# ============================================================
# SLIDE 3: Workflow Overview  (5 boxes — fits in ~12.5")
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Der KI-gestützte Feature-Workflow", 36, DARK_BG, True)

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
            "Von der Idee zum Release in 4 Phasen — jede Phase wird durch "
            "einen spezialisierten Claude Code Skill unterstützt",
            16, MEDIUM_GRAY)

flow_items = [
    ("Spezifikation", "Manuell / ChatGPT", ACCENT_PURPLE),
    ("Planung", "/plan-feature", ACCENT_BLUE),
    ("Implementierung", "/implement-feature", ACCENT_GREEN),
    ("Dokumentation", "/document-feature", ACCENT_ORANGE),
    ("Release", "GitHub Milestone", ACCENT_RED),
]
# 5×1.9 + 4×(0.2+0.3+0.2) = 9.5+2.8 = 12.3 + 0.5 start = 12.8 ✓
add_flow_boxes(slide, flow_items, top=Inches(2.2), box_w=Inches(1.9),
               box_h=Inches(1.0), start_x=Inches(0.5), gap=Inches(0.2),
               arrow_w=Inches(0.3))

desc_items = [
    ("Jede Phase hat klare Ein-/Ausgaben und ist durch GitHub Issues, "
     "Board und Milestones verknüpft.", False, DARK_GRAY),
    ("", False, DARK_GRAY),
    ("Kernprinzipien:", True, DARK_BG),
    ("   Traceability: Jedes Feature ist von User Story bis Release "
     "nachverfolgbar", False, DARK_GRAY),
    ("   Automatisierung: Repetitive Aufgaben werden durch Skills "
     "abgedeckt", False, DARK_GRAY),
    ("   Qualität: Tests (Unit, Integration, E2E) sind in den Workflow "
     "integriert", False, DARK_GRAY),
    ("   GitHub-native: Issues, Board, Milestones und Releases als "
     "Single Source of Truth", False, DARK_GRAY),
]
add_multi_text(slide, Inches(0.8), Inches(3.6), Inches(11.5), Inches(3.5),
               desc_items, 15)

# ============================================================
# SLIDE 4: Phase 1 - Spezifikation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Phase 1: Spezifikation erstellen", 36, ACCENT_PURPLE, True)

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
            "Die Grundlage jedes Features — eine durchdachte Spezifikation",
            16, MEDIUM_GRAY)

# Three approach cards: each 3.5" wide, gaps of 0.4"
# 0.8 + 3.5 + 0.4 + 3.5 + 0.4 + 3.5 = 12.1 ✓
approaches = [
    ("Manuell", ACCENT_PURPLE,
     ["Entwickler schreibt Spec", "selbständig in Markdown",
      "Volle Kontrolle über Inhalt", "Gut für einfache Features"]),
    ("ChatGPT-Dialog", ACCENT_BLUE,
     ["Iterativer Dialog mit ChatGPT", "Fragen & Antworten zur",
      "Anforderungsklärung",
      "Gut für komplexe Features"]),
    ("Hybrid", ACCENT_GREEN,
     ["Grundgerüst manuell erstellen", "Details mit KI ausarbeiten",
      "Review und Anpassung", "Empfohlener Ansatz"]),
]

for i, (title, color, lines) in enumerate(approaches):
    left = Inches(0.8) + Inches(i * 3.9)
    add_card(slide, left, Inches(2.0), Inches(3.5), Inches(2.5),
             title, lines, color)

add_textbox(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.5),
            "Struktur einer Spezifikation:", 18, DARK_BG, True)

spec_parts = [
    ("Motivation & Ziel  |  User Stories  |  Requirements  |  "
     "Architektur  |  Test-Strategie  |  Entscheidungen",
     False, DARK_GRAY),
    ("", False, DARK_GRAY),
    ("Ablageort:  docs/specs/<FEATURE-NAME>-SPEC.md", False, ACCENT_BLUE),
    ("Sprache: Deutsch mit englischen Fachbegriffen", False, MEDIUM_GRAY),
]
add_multi_text(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(2.0),
               spec_parts, 14)

# ============================================================
# SLIDE 5: Phase 2 - /plan-feature
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Phase 2: /plan-feature", 36, ACCENT_BLUE, True)

add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.6),
            "Aus der Spezifikation wird ein konkreter Implementierungsplan",
            16, MEDIUM_GRAY)

left_items = [
    ("Was der Skill macht:", True, DARK_BG),
    ("", False, DARK_GRAY),
    ("  1. Spezifikation analysieren und verstehen", False, DARK_GRAY),
    ("  2. GitHub Milestone erstellen (falls nicht vorhanden)",
     False, DARK_GRAY),
    ("  3. GitHub Issue als Epic anlegen + Board-Eintrag",
     False, DARK_GRAY),
    ("  4. Codebase analysieren: betroffene Bounded Contexts",
     False, DARK_GRAY),
    ("  5. Feature-Breakdown mit Abhängigkeiten erstellen",
     False, DARK_GRAY),
    ("  6. Implementierungsplan als PLAN-*.md erstellen",
     False, DARK_GRAY),
    ("  7. AMS UI Impact prüfen (bei Änderungen an @rsgo/core)",
     False, DARK_GRAY),
]
add_multi_text(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.5),
               left_items, 15)

# Right card: 6.8 + 5.2 = 12.0 ✓
add_card(slide, Inches(6.8), Inches(1.7), Inches(5.2), Inches(4.8),
         "Plan-Dokument enthält:",
         [
             "Ziel und Voraussetzungen",
             "Feature-Liste mit Checkboxen",
             "Neue & geänderte Dateien pro Feature",
             "Abhängigkeitskette zwischen Features",
             "API-Endpunkte (Route, Method, Permission)",
             "UI-Komponenten und State Machines",
             "Test-Strategie (Unit/Integration/E2E)",
             "Entscheidungen mit Begründung",
             "",
             "Ablage: docs/Plans/PLAN-<name>.md",
         ],
         ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.6),
            "Wichtig: Der Plan ist das zentrale Steuerungsdokument — "
            "er wird während der Implementierung laufend aktualisiert",
            14, ACCENT_BLUE, True)

# ============================================================
# SLIDE 6: Phase 3 - /implement-feature
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Phase 3: /implement-feature", 36, ACCENT_GREEN, True)

add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.6),
            "Der Plan wird Feature für Feature umgesetzt — mit Tests, "
            "PRs und Board-Updates", 16, MEDIUM_GRAY)

impl_steps = [
    ("Workflow:", True, DARK_BG),
    ("", False, DARK_GRAY),
    ("  1. Board-Status auf 'In Progress' setzen", False, DARK_GRAY),
    ("  2. Integration Branch erstellen (integration/<phase>)",
     False, DARK_GRAY),
    ("  3. Feature Branch erstellen (feature/<name>)", False, DARK_GRAY),
    ("  4. Tests schreiben (Test-First wo möglich)", False, DARK_GRAY),
    ("  5. Feature implementieren (bestehende Patterns folgen)",
     False, DARK_GRAY),
    ("  6. Build prüfen (0 Errors, 0 Warnings)", False, DARK_GRAY),
    ("  7. PR erstellen mit 'Closes #NNN'", False, DARK_GRAY),
    ("  8. Board-Status auf 'Review' setzen", False, DARK_GRAY),
    ("  9. Plan-Dokument aktualisieren (Checkboxen)", False, DARK_GRAY),
]
add_multi_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.0),
               impl_steps, 14)

# Right cards: 6.8 + 5.2 = 12.0 ✓
add_card(slide, Inches(6.8), Inches(1.6), Inches(5.2), Inches(2.2),
         "Branching-Strategie",
         [
             "main (permanent)",
             "  +-- integration/<phase> (pro Epic)",
             "       +-- feature/<name> (pro Feature)",
             "Nach Abschluss: integration -> main PR",
         ],
         ACCENT_GREEN)

add_card(slide, Inches(6.8), Inches(4.1), Inches(5.2), Inches(2.2),
         "Test-Anforderungen",
         [
             "Unit Tests (xUnit + FluentAssertions)",
             "Integration Tests (TestContainers)",
             "E2E Tests (Playwright + Screenshots)",
             "Edge Cases & Fehler-Szenarien PFLICHT!",
         ],
         ACCENT_RED)

# ============================================================
# SLIDE 7: Phase 3 continued - implement details
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "/implement-feature — Besonderheiten", 36, ACCENT_GREEN, True)

# 2 columns × 2 rows: col0 at 0.8, col1 at 0.8+6.0=6.8; width 5.5
# 6.8+5.5 = 12.3 ✓
cards_data = [
    ("Architektur-Patterns", ACCENT_GREEN, [
        "DDD: Aggregate Roots, Value Objects, Domain Events",
        "MediatR: Command/Query Separation (CQRS)",
        "FastEndpoints mit RBAC-PreProcessor",
        "EF Core + SQLite (Pre-v1.0, keine Migration)",
    ]),
    ("Board-Management", ACCENT_BLUE, [
        "Status 'In Progress' sofort beim Start",
        "Status 'Review' nach PR-Erstellung",
        "Status 'Done' NUR durch den User!",
        "Jedes Feature braucht GitHub Issue + Milestone",
    ]),
    ("AMS UI Impact", ACCENT_ORANGE, [
        "Bei Änderungen an @rsgo/core:",
        "Impact-Analyse für AMS UI durchführen",
        "Ggf. Plan im AMS Repo erstellen",
        "AMS Repo: separates privates Repository",
    ]),
    ("Qualitätssicherung", ACCENT_RED, [
        "dotnet build: 0 Errors, 0 Warnings",
        "Alle Tests müssen bestehen",
        "PR gegen Integration Branch (nicht main!)",
        "English commits, keine Footer",
    ]),
]

for i, (title, color, lines) in enumerate(cards_data):
    col = i % 2
    row = i // 2
    left = Inches(0.8) + Inches(col * 6.0)
    top = Inches(1.4) + Inches(row * 2.8)
    add_card(slide, left, top, Inches(5.5), Inches(2.4), title, lines, color)

# ============================================================
# SLIDE 8: Phase 4 - /document-feature
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Phase 4: /document-feature", 36, ACCENT_ORANGE, True)

add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.6),
            "Automatische Erstellung von E2E-Tests mit Screenshots und "
            "zweisprachiger Dokumentation", 16, MEDIUM_GRAY)

doc_steps = [
    ("Ablauf:", True, DARK_BG),
    ("", False, DARK_GRAY),
    ("  1. Feature identifizieren und UI-Workflow analysieren",
     False, DARK_GRAY),
    ("  2. Entscheiden: Erscheint Feature auf der Landing Page?",
     False, DARK_GRAY),
    ("  3. E2E-Tests mit Playwright schreiben", False, DARK_GRAY),
    ("  4. Screenshots automatisch erstellen", False, DARK_GRAY),
    ("  5. Dokumentation in Deutsch + Englisch erstellen",
     False, DARK_GRAY),
    ("  6. Optional: Landing Page Feature-Eintrag", False, DARK_GRAY),
    ("  7. Commit und PR erstellen", False, DARK_GRAY),
]
add_multi_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.0),
               doc_steps, 15)

# Right cards: 6.8 + 5.2 = 12.0 ✓
add_card(slide, Inches(6.8), Inches(1.6), Inches(5.2), Inches(2.0),
         "Screenshots",
         [
             "Playwright-basierte E2E-Tests",
             "Namenskonvention: <feature>-<##>-<desc>.png",
             "Ablage: PublicWeb/public/images/docs/",
             "Test gegen Docker Container (Port 8080)",
         ],
         ACCENT_ORANGE)

add_card(slide, Inches(6.8), Inches(3.9), Inches(5.2), Inches(2.4),
         "Dokumentation",
         [
             "Astro/Starlight Framework",
             "Zweisprachig: Deutsch + Englisch",
             "Starlight-Features: Callouts, Tabs, Tables",
             "Ablage: PublicWeb/src/content/docs/",
             "SVG-Icons für Landing Page Features",
         ],
         ACCENT_ORANGE)

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6),
            "Ergebnis: Benutzerfreundliche, bebilderte Dokumentation die "
            "direkt aus dem laufenden System generiert wird",
            14, ACCENT_ORANGE, True)

# ============================================================
# SLIDE 9: Weitere Skills
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Weitere Skills", 36, DARK_BG, True)

# 3 cards: 0.6 + 3.6 + 0.4 + 3.6 + 0.4 + 3.6 = 12.2 ✓
add_card(slide, Inches(0.6), Inches(1.4), Inches(3.6), Inches(5.0),
         "/fix-bug",
         [
             "Red-Green Test-Ansatz:",
             "",
             "1. Bug-Issue von GitHub lesen",
             "2. Board -> 'In Progress'",
             "3. Bugfix Branch: bugfix/<name>",
             "4. RED: Test schreiben der fehlschlägt",
             "5. GREEN: Minimaler Fix",
             "6. Build + alle Tests prüfen",
             "7. PR mit 'Fixes #NNN'",
             "8. Board -> 'Review'",
             "",
             "Kein Refactoring, nur der Fix!",
         ],
         ACCENT_RED)

add_card(slide, Inches(4.6), Inches(1.4), Inches(3.6), Inches(5.0),
         "/report-bug",
         [
             "Bug melden ohne zu fixen:",
             "",
             "1. Bug-Beschreibung verstehen",
             "2. Code-Analyse & Docker Logs",
             "3. GitHub Issue erstellen mit:",
             "   - Description",
             "   - Expected Behavior",
             "   - Reproduction Steps",
             "   - Environment",
             "   - Analysis",
             "4. Issue auf Board -> 'Backlog'",
             "",
             "Nächster Schritt: /fix-bug",
         ],
         ACCENT_ORANGE)

add_card(slide, Inches(8.6), Inches(1.4), Inches(3.6), Inches(5.0),
         "/docker-dev",
         [
             "Container für manuelles Testen:",
             "",
             "1. Clean Mode? (Volumes löschen?)",
             "2. docker compose down [-v]",
             "3. docker compose build",
             "4. docker compose up -d",
             "5. Status prüfen",
             "",
             "URL: http://localhost:8080",
             "",
             "Volumes löschen = Setup Wizard",
             "Volumes behalten = Bestandsdaten",
         ],
         ACCENT_BLUE)

# ============================================================
# SLIDE 10: GitHub Integration
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "GitHub-Integration", 36, DARK_BG, True)

add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.6),
            "GitHub als Single Source of Truth — Issues, Board, Milestones "
            "und Releases sind vollständig integriert", 16, MEDIUM_GRAY)

# 3 cards top row: 0.6 + 3.6 + 0.4 + 3.6 + 0.4 + 3.6 = 12.2 ✓
add_card(slide, Inches(0.6), Inches(1.9), Inches(3.6), Inches(2.3),
         "GitHub Issues",
         [
             "Jedes Feature = Epic Issue",
             "Labels: epic, feature, bug, docs",
             "Milestone-Zuweisung (Pflicht!)",
             "Sub-Issues für einzelne Features",
             "Auto-Close via 'Closes #NNN' in PR",
         ],
         ACCENT_BLUE)

add_card(slide, Inches(4.6), Inches(1.9), Inches(3.6), Inches(2.3),
         "Project Board",
         [
             "Status-Workflow:",
             "  Backlog -> In Progress -> Review -> Done",
             "Skills setzen Status automatisch",
             "'Done' nur durch User (Qualitäts-Gate)",
             "Zentrale Übersicht aller Arbeit",
         ],
         ACCENT_GREEN)

add_card(slide, Inches(8.6), Inches(1.9), Inches(3.6), Inches(2.3),
         "Milestones & Releases",
         [
             "1 Milestone = 1 Release-Version",
             "Milestone schließen = Release!",
             "Automatische Release Notes",
             "Automatischer Docker Build",
             "Versionierung: v0.XX.0",
         ],
         ACCENT_PURPLE)

# 2 cards bottom row: 0.6 + 5.4 + 0.4 + 5.4 = 11.8 ✓
add_card(slide, Inches(0.6), Inches(4.5), Inches(5.4), Inches(2.4),
         "CI/CD Pipeline",
         [
             "Pull Request -> CI automatisch:",
             "  - dotnet restore, build, test",
             "  - pnpm install, lint, type-check, build",
             "  - Test-Ergebnisse als Artifacts (7 Tage)",
             "",
             "Milestone Close -> Release automatisch:",
             "  - Release Notes generieren",
             "  - GitHub Release + Tag erstellen",
             "  - Docker Build triggern",
         ],
         DARK_GRAY)

add_card(slide, Inches(6.4), Inches(4.5), Inches(5.8), Inches(2.4),
         "Release Drafter (Automatische Kategorisierung)",
         [
             "Features, Bug Fixes, Security, Docs, Maintenance",
             "Auto-Labels basierend auf:",
             "  - Branch-Namen (fix/*, feature/*, ...)",
             "  - Geänderte Dateien (*.md, *.csproj, ...)",
             "",
             "Versionierung automatisch:",
             "  major Label -> Major, feature -> Minor",
             "  bug/fix/docs -> Patch (Standard)",
         ],
         DARK_GRAY)

# ============================================================
# SLIDE 11: GitHub Flow Diagram  (8 boxes — tighter layout)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "GitHub-Workflow im Detail", 36, DARK_BG, True)

# 8 boxes: 8×1.25 + 7×(0.08+0.15+0.08) = 10.0+2.17 = 12.17 + 0.4 = 12.57 ✓
flow_data = [
    ("Issue erstellen", "/plan-feature", ACCENT_PURPLE),
    ("Board: Backlog", "Automatisch", LIGHT_GRAY),
    ("Board:\nIn Progress", "/implement", ACCENT_BLUE),
    ("PR erstellen", "Closes #NNN", ACCENT_GREEN),
    ("CI Pipeline", "Automatisch", ACCENT_ORANGE),
    ("Board: Review", "Automatisch", ACCENT_ORANGE),
    ("PR Merge", "User Review", ACCENT_GREEN),
    ("Board: Done", "User setzt", ACCENT_RED),
]
add_flow_boxes(slide, flow_data, top=Inches(1.5), box_w=Inches(1.25),
               box_h=Inches(1.1), start_x=Inches(0.4), gap=Inches(0.08),
               arrow_w=Inches(0.15))

add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.5),
            "Release-Prozess:", 20, DARK_BG, True)

# 5 boxes: 5×2.0 + 4×(0.1+0.25+0.1) = 10.0+1.8 = 11.8 + 0.5 = 12.3 ✓
release_flow = [
    ("Alle Features done", "", ACCENT_GREEN),
    ("Integration -> main", "Finaler PR", ACCENT_BLUE),
    ("Milestone schließen", "User Aktion", ACCENT_PURPLE),
    ("Release erstellt", "Automatisch", ACCENT_ORANGE),
    ("Docker Build", "Automatisch", ACCENT_RED),
]
add_flow_boxes(slide, release_flow, top=Inches(3.6), box_w=Inches(2.0),
               box_h=Inches(0.9), start_x=Inches(0.5), gap=Inches(0.1),
               arrow_w=Inches(0.25))

key_items = [
    ("", False, DARK_GRAY),
    ("Vorteile dieser Integration:", True, DARK_BG),
    ("  Lückenlose Nachverfolgbarkeit: Issue -> PR -> Release",
     False, DARK_GRAY),
    ("  Automatische Status-Updates durch die Skills", False, DARK_GRAY),
    ("  User behält Kontrolle über finale Freigaben "
     "(Done, Merge, Release)", False, DARK_GRAY),
    ("  CI/CD stellt Qualität bei jedem PR sicher", False, DARK_GRAY),
]
add_multi_text(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.2),
               key_items, 14)

# ============================================================
# SLIDE 12: Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Zusammenfassung", 36, WHITE, True)

skills = [
    ("/plan-feature",
     "Spezifikation -> Plan + GitHub Issue + Board", ACCENT_BLUE),
    ("/implement-feature",
     "Plan -> Code + Tests + PR + Board-Updates", ACCENT_GREEN),
    ("/document-feature",
     "Feature -> E2E Tests + Screenshots + Docs (DE/EN)", ACCENT_ORANGE),
    ("/fix-bug",
     "Bug-Issue -> Red-Green Fix + PR", ACCENT_RED),
    ("/report-bug",
     "Bug -> GitHub Issue + Board-Eintrag", ACCENT_ORANGE),
    ("/docker-dev",
     "Container bauen + starten für manuelles Testen", ACCENT_BLUE),
]

for i, (skill, desc, color) in enumerate(skills):
    y = Inches(1.4) + Inches(i * 0.7)
    add_textbox(slide, Inches(1.0), y, Inches(3.0), Inches(0.5), skill, 18,
                color, True, font_name="Consolas")
    add_textbox(slide, Inches(4.2), y, Inches(8), Inches(0.5), desc, 16,
                LIGHT_GRAY)

add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2),
            "Der KI-gestützte Workflow ermöglicht es, Features schnell und "
            "konsistent von der Idee bis zum\nRelease umzusetzen — mit "
            "integrierter Qualitätssicherung und lückenloser "
            "Nachverfolgbarkeit.",
            16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = r"c:\proj\ReadyStackGo\docs\KI-Feature-Workflow-ReadyStackGo.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
